"""VERO Parser — PPTX: Slide-aware parsing with table and speaker notes support.

Extracts content from PowerPoint presentations:
  - Each slide becomes a ## Slide N section
  - Title shapes become ### headings
  - Tables → Markdown table format
  - Speaker notes → blockquoted notes
  - Images → counted for metadata
"""

import logging
from pathlib import Path

from pptx import Presentation

from app.parsers.contracts import ParsedDocument, table_to_markdown

logger = logging.getLogger(__name__)


async def parse_pptx(filepath: str) -> dict:
    """Parse a PPTX file with slide structure, tables, and speaker notes.

    Returns {"text": str, "metadata": dict, "parsed_doc": ParsedDocument}
    """
    filename = Path(filepath).name
    prs = Presentation(filepath)

    slides_md = []
    tables_raw = []
    image_count = 0

    for i, slide in enumerate(prs.slides):
        slide_parts = [f"## Slide {i + 1}"]

        for shape in slide.shapes:
            # Count images (shape_type 13 = Picture)
            if shape.shape_type == 13:
                image_count += 1
                continue

            if not shape.has_text_frame and not shape.has_table:
                continue

            # Extract text with title detection
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # Title shapes get heading treatment
                    if shape.name and shape.name.lower().startswith("title"):
                        slide_parts.append(f"### {text}")
                    else:
                        slide_parts.append(text)

            # Extract tables
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows and len(rows) >= 2:
                    tables_raw.append(rows)
                    md = table_to_markdown(rows)
                    if md:
                        slide_parts.append(f"\n{md}")

        # Extract speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"\n> **Notes:** {notes}")

        slides_md.append("\n\n".join(slide_parts))

    markdown_text = "\n\n---\n\n".join(slides_md)

    parsed_doc = ParsedDocument(
        source_type="pptx",
        filename=filename,
        markdown_text=markdown_text,
        page_count=len(prs.slides),
        has_images=image_count > 0,
        has_tables=len(tables_raw) > 0,
        has_equations=False,
        is_slide_format=True,
        complexity="slide_mode",
        image_descriptions=[],
        tables_raw=tables_raw,
        metadata={
            "slide_count": len(prs.slides),
            "table_count": len(tables_raw),
            "image_count": image_count,
        },
    )

    logger.info(
        "PPTX parsed [%s]: %d chars, %d slides, %d tables, %d images",
        filename, len(markdown_text), len(prs.slides), len(tables_raw), image_count,
    )

    return {
        "text": markdown_text,
        "metadata": parsed_doc.metadata,
        "parsed_doc": parsed_doc,
    }
